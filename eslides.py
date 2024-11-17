from pptx import Presentation

# Criação da apresentação
presentation = Presentation()

# Slide 1: Título
slide = presentation.slides.add_slide(presentation.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Caminho Mínimo de Dijkstra"
subtitle.text = "Seminário sobre conceitos, características e aplicações"

# Slide 2: Introdução
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Introdução"
content.text = (
    "- O problema do caminho mínimo consiste em encontrar o menor custo entre dois nós de um grafo.\n"
    "- Utilizado em diversas áreas como navegação GPS, redes de computadores e logística.\n"
    "- Baseado no algoritmo criado por Edsger W. Dijkstra em 1956."
)

# Slide 3: Principais Conceitos
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Principais Conceitos"
content.text = (
    "- Grafos: Conjunto de vértices e arestas com pesos associados.\n"
    "- Caminho mínimo: Soma dos menores pesos entre nós do grafo.\n"
    "- Algoritmo de Dijkstra: Método guloso para resolver o problema."
)

# Slide 4: Características
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Características"
content.text = (
    "- Funciona apenas com pesos não negativos.\n"
    "- Complexidade computacional:\n"
    "  - Com heap: O((V + E) log V)\n"
    "  - Sem heap: O(V²)\n"
    "- É um algoritmo guloso, escolhendo sempre o menor custo conhecido."
)

# Slide 5: Aplicações
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Aplicações"
content.text = (
    "- Navegação GPS: Encontra rotas otimizadas.\n"
    "- Redes de computadores: Protocolos como OSPF utilizam Dijkstra.\n"
    "- Planejamento de transporte e logística.\n"
    "- Caminhos em jogos e tabuleiros."
)

# Slide 6: Funcionamento do Algoritmo
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Funcionamento do Algoritmo"
content.text = (
    "1. Inicialize as distâncias: Origem = 0, outros = infinito.\n"
    "2. Escolha o nó com menor custo acumulado.\n"
    "3. Atualize distâncias dos vizinhos se encontrar caminhos mais curtos.\n"
    "4. Repita até visitar todos os nós ou alcançar o destino."
)

# Slide 7: Extensões e Comparações
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Extensões e Comparações"
content.text = (
    "- Algoritmo de Bellman-Ford: Funciona com pesos negativos.\n"
    "- Algoritmo A*: Utiliza heurísticas para otimizar buscas.\n"
    "- Floyd-Warshall: Calcula caminhos mínimos entre todos os pares de vértices."
)

# Slide 8: Conclusão
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Conclusão"
content.text = (
    "- O algoritmo de Dijkstra é eficiente e amplamente usado em problemas reais.\n"
    "- Suas aplicações abrangem diversas áreas tecnológicas.\n"
    "- Limitações como pesos negativos exigem algoritmos complementares."
)

# Salvar apresentação
file_path = "Caminho_Minimo_Dijkstra_Seminario_Corrigido.pptx"
presentation.save(file_path)
print(f"Apresentação salva como {file_path}")